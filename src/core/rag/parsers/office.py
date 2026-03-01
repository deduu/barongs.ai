from __future__ import annotations

import asyncio
import io

try:
    from pptx import Presentation  # type: ignore[import-untyped]
except ImportError:
    Presentation = None  # type: ignore[assignment,misc]

try:
    from openpyxl import load_workbook  # type: ignore[import-untyped]
except ImportError:
    load_workbook = None  # type: ignore[assignment]

_PPTX_MSG = "PPTX parsing requires python-pptx. Install with: pip install barongsai[rag-parsers]"
_XLSX_MSG = "Excel parsing requires openpyxl. Install with: pip install barongsai[rag-parsers]"


class PowerPointParser:
    """Parser for PowerPoint (.pptx) files."""

    @property
    def supported_extensions(self) -> frozenset[str]:
        return frozenset({".pptx"})

    async def parse(self, raw: bytes, filename: str) -> str:
        if Presentation is None:
            raise ImportError(_PPTX_MSG)
        return await asyncio.to_thread(self._extract, raw)

    def _extract(self, raw: bytes) -> str:
        prs = Presentation(io.BytesIO(raw))
        slides: list[str] = []
        for i, slide in enumerate(prs.slides, 1):
            texts: list[str] = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    texts.append(shape.text_frame.text)
            slides.append(f"[Slide {i}]\n" + "\n".join(texts))
        return "\n\n".join(slides)


class ExcelParser:
    """Parser for Excel (.xlsx) files."""

    @property
    def supported_extensions(self) -> frozenset[str]:
        return frozenset({".xlsx"})

    async def parse(self, raw: bytes, filename: str) -> str:
        if load_workbook is None:
            raise ImportError(_XLSX_MSG)
        return await asyncio.to_thread(self._extract, raw)

    def _extract(self, raw: bytes) -> str:
        wb = load_workbook(io.BytesIO(raw), read_only=True, data_only=True)
        sheets: list[str] = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows: list[str] = []
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) if c is not None else "" for c in row]
                rows.append("\t".join(cells))
            sheets.append(f"[Sheet: {sheet_name}]\n" + "\n".join(rows))
        wb.close()
        return "\n\n".join(sheets)

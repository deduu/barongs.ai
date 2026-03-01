from __future__ import annotations

import io

import pytest
from openpyxl import Workbook
from pptx import Presentation
from pptx.util import Inches

from src.core.rag.parsers.office import ExcelParser, PowerPointParser


def _make_pptx(slide_texts: list[str]) -> bytes:
    """Create a minimal PPTX with one text box per slide."""
    prs = Presentation()
    for text in slide_texts:
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
        text_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
        text_box.text_frame.text = text
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _make_xlsx(sheets: dict[str, list[list[str]]]) -> bytes:
    """Create a minimal XLSX with named sheets and row data."""
    wb = Workbook()
    # Remove default sheet
    wb.remove(wb.active)
    for name, rows in sheets.items():
        ws = wb.create_sheet(title=name)
        for row in rows:
            ws.append(row)
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


class TestPowerPointParser:
    def setup_method(self) -> None:
        self.parser = PowerPointParser()

    def test_supported_extensions(self) -> None:
        assert self.parser.supported_extensions == frozenset({".pptx"})

    @pytest.mark.asyncio
    async def test_parse_extracts_slide_text(self) -> None:
        pptx_bytes = _make_pptx(["First slide content", "Second slide content"])
        result = await self.parser.parse(pptx_bytes, "test.pptx")
        assert "First slide content" in result
        assert "Second slide content" in result
        assert "[Slide 1]" in result
        assert "[Slide 2]" in result

    @pytest.mark.asyncio
    async def test_parse_empty_presentation(self) -> None:
        pptx_bytes = _make_pptx([])
        result = await self.parser.parse(pptx_bytes, "empty.pptx")
        assert result.strip() == ""


class TestExcelParser:
    def setup_method(self) -> None:
        self.parser = ExcelParser()

    def test_supported_extensions(self) -> None:
        exts = self.parser.supported_extensions
        assert ".xlsx" in exts

    @pytest.mark.asyncio
    async def test_parse_extracts_sheet_data(self) -> None:
        xlsx_bytes = _make_xlsx({
            "Sales": [["Product", "Revenue"], ["Widget", "1000"]],
            "Costs": [["Item", "Amount"], ["Labor", "500"]],
        })
        result = await self.parser.parse(xlsx_bytes, "test.xlsx")
        assert "[Sheet: Sales]" in result
        assert "Widget" in result
        assert "1000" in result
        assert "[Sheet: Costs]" in result
        assert "Labor" in result
